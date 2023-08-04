
from plotly.offline import plot
import plotly.graph_objs as go
import plotly.express as px
import pandas as pd

def scatter(dflastdate):
    # df.dropna(subset='continent',inplace=True)
    # dflastdate=df[(df['date'].str[-2:]=='07')]

    dflastdate=dflastdate.groupby(['continent','date']).sum().reset_index()
    scatter = px.scatter(dflastdate, x="date", y="total_cases", color="continent",
                  hover_data=['total_deaths','total_vaccinations'], size='total_deaths')
    scatter.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        margin = dict(t=26, l=5, r=5, b=5),
        font_family="Comic Sans MS",
        font_color="white",
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    scatter.write_image("static/index_scatter.png")
    plt_div=plot(scatter,output_type='div')
    return plt_div


def histogram1(df):
    dflastdate=df[(df['date'].str[:]=='2023-06-07')]
    histogram=px.bar(dflastdate, x="continent",y="total_cases",color='location')
    histogram.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        font_family="Comic Sans MS",
        font_color="white",
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    plt_div= plot(histogram, output_type='div')
    histogram.write_image("static/index_hist1.png")
    return plt_div

def bargraphtop5(df):
    # df.dropna(subset='continent',inplace=True)
    dflastdate=df[df['date']=='2023-06-07']
    top5countrydata=pd.DataFrame()
    for cont in dflastdate.continent.unique():
        dftmp=dflastdate[dflastdate['continent']==cont]
        dftmp=dftmp.sort_values(ascending=False,by=['total_cases']).head(5)
        top5countrydata=top5countrydata.append(dftmp)
    barchart=px.bar(top5countrydata,y="continent",x="total_cases",color="location",orientation='h',hover_data=['total_cases','total_deaths'],barmode='group')
    barchart.update_traces(width=.7)
    barchart.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        font_family="Comic Sans MS",
        font_color="white",
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    barchart.write_image("static/index_bargraphtop5.png")
    plt_div= plot(barchart, output_type='div')
    return plt_div

def sunburst1(df):
    # df.dropna(subset='continent',inplace=True)
    dflastdate=df[df['date']=='2023-06-07']
    sun = px.sunburst(dflastdate, path=['continent', 'location'], values='total_cases',title="TOTAL CASES",color_discrete_sequence = ['blue'])
    sun.update_traces(textinfo="label+percent parent")
    sun.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        margin = dict(t=26, l=5, r=5, b=5),
        font_family="Comic Sans MS",
        font_color="white",
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    sun.write_image("static/index_sunburst1.png")
    plt_div=plot(sun,output_type='div')
    return plt_div

def sunburst2(df):
    # df.dropna(subset='continent',inplace=True)
    dflastdate=df[df['date']=='2023-06-07']
    sun = px.sunburst(dflastdate, 
                    path=['continent', 'location'], 
                    values='total_deaths',
                    title="TOTAL DEATHS",
                    color_discrete_sequence = ['red'])
    sun.update_traces(textinfo="label+percent parent")
    sun.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        margin = dict(t=26, l=5, r=5, b=5),
        font_family="Comic Sans MS",
        font_color="white",
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    sun.write_image("static/index_sunburst2.png")
    plt_div=plot(sun,output_type='div')
    return plt_div

def sunburst3(df):
    # df.dropna(subset='continent',inplace=True)
    dflastdate=df[df['date']=='2023-06-07']
    sun = px.sunburst(dflastdate, 
                    path=['continent', 'location'], 
                    values='population',
                    title="TOTAL POPULATION",
                    color_discrete_sequence = ['green'])
    sun.update_traces(textinfo="label+percent parent")
    sun.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        margin = dict(t=26, l=5, r=5, b=5),
        font_family="Comic Sans MS",
        font_color="white",
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    sun.write_image("static/index_sunburst3.png")
    plt_div=plot(sun,output_type='div')
    return plt_div

def line1(df):
    # df.dropna(subset='continent',inplace=True)
    dflastdate=df[(df['date'].str[-2:]=='07')]

    dflastdate=dflastdate.groupby(['continent','date']).sum().reset_index()
    line = px.line(dflastdate, x="date", y="new_cases", color='continent')
    line.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        margin = dict(t=26, l=5, r=5, b=5),
        font_family="Comic Sans MS",
        font_color="white",
        height=250,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    line.write_image("static/index_line1.png")
    plt_div=plot(line,output_type='div')
    return plt_div

def line2(df):
    # df.dropna(subset='continent',inplace=True)
    dflastdate=df[(df['date'].str[-2:]=='07')]

    dflastdate=dflastdate.groupby(['continent','date']).sum().reset_index()
    line = px.line(dflastdate, x="date", y="positive_rate", color='continent')
    line.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        margin = dict(t=26, l=5, r=5, b=5),
        font_family="Comic Sans MS",
        font_color="white",
        height=250,
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    line.write_image("static/index_line2.png")
    plt_div=plot(line,output_type='div')
    return plt_div
    



    #country graph
def con_line1(dfind):
    # df.dropna(subset='continent',inplace=True)
    # dfind=df[(df['location'].str.lower()==sel_con.lower()) & (df['date'].str[-2:]=='01')]

    line = px.line(dfind, x="date", y="total_cases")
    line.update_layout({
            'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
            'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
            },
            margin = dict(t=26, l=5, r=5, b=5),
            font_family="Comic Sans MS",
            font_color="white",
            height=250,
            xaxis=dict(showgrid=False),
            yaxis=dict(showgrid=False))
    line.update_traces(line_color='pink')
    plt_div=plot(line,output_type='div')
    line.write_image("static/country_line1.png")
    return plt_div

def con_line2(dfind):
    # df.dropna(subset='continent',inplace=True)
    # dfind=df[(df['location'].str.lower()==sel_con.lower()) & (df['date'].str[-2:]=='01')]

    line = px.line(dfind, x="date", y="new_cases")
    line.update_layout({
            'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
            'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
            },
            margin = dict(t=26, l=5, r=5, b=5),
            font_family="Comic Sans MS",
            font_color="white",
            height=250,
            xaxis=dict(showgrid=False),
            yaxis=dict(showgrid=False))
    line.update_traces(line_color='yellow')
    plt_div=plot(line,output_type='div')
    line.write_image("static/country_line2.png")
    return plt_div


def con_line3(dfind):
    # df.dropna(subset='continent',inplace=True)
    # dfind=df[(df['location'].str.lower()==sel_con.lower()) & (df['date'].str[-2:]=='07')]

    line = px.line(dfind, x="date", y="total_vaccinations")
    line.update_layout({
            'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
            'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
            },
            margin = dict(t=26, l=5, r=5, b=5),
            font_family="Comic Sans MS",
            font_color="white",
            height=250,
            xaxis=dict(showgrid=False),
            yaxis=dict(showgrid=False))
    line.update_traces(line_color='green')
    plt_div=plot(line,output_type='div')
    line.write_image("static/country_line3.png")
    return plt_div

def con_line4(dfind):
    # df.dropna(subset='continent',inplace=True)
    # dfind=df[(df['location'].str.lower()==sel_con.lower()) & (df['date'].str[-2:]=='01')]

    line = px.line(dfind, x="date", y="total_deaths")
    line.update_layout({
            'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
            'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
            },
            
            margin = dict(t=26, l=5, r=5, b=5),
            font_family="Comic Sans MS",
            font_color="white",
            height=250,
            xaxis=dict(showgrid=False),
            yaxis=dict(showgrid=False))
    line.update_traces(line_color='red')
    plt_div=plot(line,output_type='div')
    line.write_image("static/country_line4.png")
    return plt_div


def sunburst_con(df,sel_con):
    df.dropna(subset='continent',inplace=True)
    selected_continent=df[df['location'].str.lower()==sel_con.lower()]['continent'].iloc[0]
    dflastdate=df[(df['continent']==selected_continent) & (df['date']=='2023-06-07')]
    sun = px.sunburst(dflastdate, 
                    path=['continent', 'location'], 
                    values='population',
                    title="Population",

                    color_discrete_sequence = ['blue'])
    sun.update_traces(textinfo="label+percent parent")
    sun.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        margin = dict(t=26, l=5, r=5, b=5),
        font_family="Comic Sans MS",
        font_color="white",
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    plt_div=plot(sun,output_type='div')
    sun.write_image("static/country_sunburst1.png")
    return plt_div



def bargraphtop(df,sel_con):
    df.dropna(subset='continent',inplace=True)
    dflastdate=df[df['date'].str[-5:]=='12-31']
    dflastdate2=df[df['date']=='2023-03-01']
    my_current_country=sel_con
    top5_countries_list=['United States','China','India','France','Germany']
    top_5_country_data=dflastdate[dflastdate['location'].isin(top5_countries_list)]
    top_5_country_data=top_5_country_data.append(dflastdate2[dflastdate2['location'].isin(top5_countries_list)])
    top_5_country_data['year']=top_5_country_data['date'].str[:4]
    my_country_data=dflastdate2[dflastdate2['location']==my_current_country]
    my_country_data=my_country_data.append(dflastdate[dflastdate['location']==my_current_country])
    my_country_data['year']=my_country_data['date'].str[:4]
    my_country_data=my_country_data.sort_values(by=['total_cases'])
    fig=px.bar(top_5_country_data[['location','total_cases','year']], x="year", y="total_cases",
                color="location", hover_data=['total_cases'],
                barmode = 'group')
    fig.add_scatter(x=my_country_data["year"],y=my_country_data["total_cases"],name=my_current_country)
   


    fig.update_layout({
        'plot_bgcolor': 'rgba(0, 0, 0, 0.300)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0.300)',
        },
        font_family="Comic Sans MS",
        title='COMPARISON',
        title_x=0.5,
        font_color="white",
        xaxis=dict(showgrid=False),
        yaxis=dict(showgrid=False))
    plt_div= plot(fig, output_type='div')
    fig.write_image("static/country_barline.png")
    return plt_div


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image


def download_ppt(name_on_first_slide,sel_con,continent):
        prs = Presentation('../covid-19 ppt.pptx')
        #Editting Name on First Slide
        prs.slides[0].shapes[-1].text='- '+name_on_first_slide.upper()
        prs.slides[0].shapes[-1].text_frame.paragraphs[0].font.size = Pt(22)
        prs.slides[0].shapes[-1].text_frame.paragraphs[0].font.name = 'Arial'
        prs.slides[0].shapes[-1].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        #country name replacement
        # Specify the text to find and replace 
        find_text = 'India' 
        replace_text =sel_con

        def find_replace_text(presentation, find_text, replace_text): 
            for slide in presentation.slides: 
                for shape in slide.shapes: 
                    if shape.has_text_frame: 
                        for paragraph in shape.text_frame.paragraphs: 
                            for run in paragraph.runs: 
                                if find_text in run.text: 
                                    run.text = run.text.replace(find_text, replace_text) 
                                    
                                    
        # Call the function to find and replace the text 
        find_replace_text(prs, find_text, replace_text) 

        #continent name replacement
        # Specify the text to find and replace 
        find_text_continent = 'Asia' 
        replace_text_continent =continent

                                    
                                    
        # Call the function to find and replace the text 
        find_replace_text(prs, find_text_continent, replace_text_continent) 


        #Slide 12 Image creation
        def insert_pic(slideno,left,top,width,height,imgpath):
            slide = prs.slides[slideno]
            pic=slide.shapes.add_picture(imgpath, left, top,
                                        width, height) 
        def insert_text(slide_no,shape_no,newtext,fsize,fname):
            prs.slides[slide_no].shapes[shape_no].text=newtext
            prs.slides[slide_no].shapes[shape_no].text_frame.paragraphs[0].font.size = Pt(fsize)
            prs.slides[slide_no].shapes[shape_no].text_frame.paragraphs[0].font.name = fname
            prs.slides[slide_no].shapes[shape_no].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        #update the heading
        # Slide 12
        # COVID-19 Cases Analysis in India: Total Cases vs. New Cases
        # for i,pl in enumerate(prs.slides[11].shapes):
        #     print(i,pl.text)
        # insert_text(11,2,'COVID-19 Cases Analysis in '+sel_con+': Total Cases vs. New Cases',28,'Modern Love (bold)')
            

        # insert_text(11,4,'Engage in a detailed analysis of '+sel_con+' battle against COVID-19 with the illuminating Line Charts.'+
        # 'Uncover the trajectory of the pandemic within ' +sel_con+', showcased through two insightful line charts: "Total Cases" and New Cases.'+
        # 'The "Total Cases" chart presents '+sel_con+' overall case count, providing a panoramic view of the pandemic scale and progression within the nation.'+
        # 'Observe how the case count evolved over time, revealing critical milestones and challenges faced during the pandemic course.'+
        # 'Delve into the intricacies of the pandemic spread with the "New Cases" chart, depicting the daily or weekly influx of new infections.'+
        # 'Identify periods of surges, dips, or stability, providing valuable insights into the effectiveness of containment measures and vaccination campaigns.'+
        # 'Compare and contrast the trends between the two charts to uncover connections between the accumulation of cases and the fluctuations in new infections. ',16,'Arial')


        insert_pic(11,Inches(0),Inches(0),Inches(6.49),Inches(3.58),'../heart_analysis/static/country_line1.png')
        insert_pic(11,Inches(6.55),Inches(0),Inches(6.70),Inches(3.58),'../heart_analysis/static/country_line2.png')

        # #update the heading
        # # Slide 13
        # # COVID-19 Impact: Vaccination and Total Deaths in India
        # insert_text(12,2,'COVID-19 Impact: Vaccination and Total Deaths in '+sel_con,28,'Modern Love (bold)')

        # insert_text(12,4,'Unveil the pivotal aspects of  ' +sel_con+' COVID-19 journey with two crucial line charts: "Vaccination" and "Total Deaths."'+
        # 'Gain profound insights into ' +sel_con+' vaccination progress and its impact on the nation resilience in the face of the pandemic.'+
        # 'The "Vaccination" chart showcases the cumulative doses administered, reflecting ' +sel_con+' strides towards herd immunity and protection.'+
        # 'Observe the acceleration in vaccinations, indicating a concerted effort in safeguarding the population.'+
        # 'Delve into the effectiveness of vaccination campaigns, potentially curbing infection rates and reducing severe outcomes.'+
        # 'The "Total Deaths" chart portrays the solemn toll of the pandemic on ' +sel_con+', capturing the lives lost to COVID-19.'+
        # 'Pay tribute to those who lost their lives, understanding the human cost of the pandemic impact.'+
        # 'Analyze the intersection of vaccination efforts and the impact on reducing fatalities.'
        #  ,14,'Arial')


        insert_pic(12,Inches(0),Inches(0),Inches(6.49),Inches(3.58),'../heart_analysis/static/country_line3.png')
        insert_pic(12,Inches(6.55),Inches(0),Inches(6.70),Inches(3.58),'../heart_analysis/static/country_line4.png')


        # # Slide 14
        # # Pandemic Perspectives: Top 5 Countries vs. India
        # # left 5.85
        # # width 7.36
        # # height 5.86
        # # top 0.85
        # for i,pl in enumerate(prs.slides[13].shapes):
        #     print(i,pl.text)

        # insert_text(12,2,'COVID-19 Impact: Vaccination and Total Deaths in '+sel_con,28,'Modern Love (bold)')


        insert_pic(13,Inches(5.85),Inches(0.85),Inches(7.36),Inches(5.86),'../heart_analysis/static/country_barline.png')



        # # Slide 15
        # #  Continental Insights: Sunburst Chart of Population in Asia​

        # # Embark on an enlightening journey into the heart of Asia with the mesmerizing "Sunburst Chart."​

        # # Unveil the total population of Asia, the continent to which India belongs, offering a holistic understanding of the region's significance in the global fight against COVID-19.​

        # # The "Sunburst Chart" artfully visualizes the population distribution across different tiers, providing a comprehensive view of the continent's demographics.​

        # # Explore the hierarchy of population segments, from countries to regions, cities, and beyond, to comprehend Asia's immense diversity.​

        # # Understand how India, as a prominent country within Asia, fits into this broader tapestry, recognizing its role within the continental context.​

        # # Gain valuable insights into Asia's capacity to mobilize resources, expertise, and collective responses in combatting the pandemic.​

        # # Empower policymakers and the audience with nuanced knowledge, guiding informed strategies for regional cooperation and support.​

        # # ​left 5.85
        # # width 7.36
        # # height 5.86
        # # top 0.85
        #insert_text(14,6,' Continental Insights: Sunburst Chart of Population in '+continent,28,'Modern Love (bold)')


        insert_pic(14,Inches(5.85),Inches(0.85),Inches(7.45),Inches(5.75),'../heart_analysis/static/country_sunburst1.png')

        prs.save('../heart_analysis/static/covid-19-editied.pptx')
        



